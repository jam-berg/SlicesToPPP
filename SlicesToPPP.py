#!/usr/bin/env python3
import collections.abc
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Cm
import os
import glob

base_path = os.environ['HOME'] + '/ETHZ/Semester 3/Projekt&Praktika III/C8 Radikalische Polymerisation/STL Examples/QuadraTresHelicesRndCrnr/'
img_path = 'QuadraTresHelicesRndCrnr00120.png'
filepath_list = sorted(glob.glob(base_path + '/*.png'))
print(filepath_list)
	

prs = Presentation()
slide_masters = prs.slide_masters
masterfill = slide_masters[0].background.fill
masterfill.solid()
masterfill.fore_color.rgb = RGBColor(0, 0, 0)
title_slide_layout = prs.slide_layouts[9]

for filepath in filepath_list:
	slide = prs.slides.add_slide(title_slide_layout)
	pic = slide.shapes.add_picture(filepath, Cm(7.34), 0, height=prs.slide_height)
	slide = prs.slides.add_slide(title_slide_layout)

prs.save('test.pptx')

